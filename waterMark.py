from os import listdir
from wand.image import Image
from pptx import Presentation
from pptx.util import Inches

image_destination = input("Enter the image destination: ") + "\\"
watermark_name = input("Enter the watermark's name:")
watermark_file = image_destination + watermark_name + ".png"
imagesList = listdir(image_destination)
path = input("Enter the output destination: ") + "\\"


def get_all_images(image_path, watermark, output_path):
    imagesList = listdir(image_path)

    for image in imagesList:
        img = (image_path + image)
        image_check = str(img).split(".")
        if image_check[1] == 'jpg':
            with Image(filename=img) as pic:
                with Image(filename=watermark) as wm:
                    wm.transparentize(0.11)
                    with Image(width=pic.width, height=pic.height) as final:
                        final.composite(image=pic, left=0, top=0)
                        final.composite(image=wm, left=1, top=1)
                        name = image.replace(".jpg", "")
                        final.save(filename=output_path + name + ".png")


def presentation(images_path):
    prs = Presentation()

    lyt = prs.slide_layouts[1] 

    # images from the path
    imagesList = listdir(images_path)

    for image in imagesList:
        slide = prs.slides.add_slide(lyt)  
        prs.slide_width = Inches(56)
        prs.slide_height = Inches(56)
        title = slide.shapes.title  
        subtitle = slide.placeholders[1]  
        name = image.replace(".jpg.png", "")
        title.text = f"{name}"  
        subtitle.text = f"{name}" 

        img = (images_path + image)

        slide.shapes.add_picture(img, left=Inches(1), top=Inches(2.5))

    prs.save(f'{images_path}\output.pptx')  

get_all_images(image_destination, watermark_file, path)
presentation(path)
